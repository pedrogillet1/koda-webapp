#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract text from PowerPoint (.pptx) files
Supports text extraction from slides, tables, notes, and nested shapes
"""

import sys
import json
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set UTF-8 encoding for stdout on Windows
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

def extract_text_from_slide(slide, slide_number):
    """Extract all text from a single slide including tables and notes"""
    texts = []

    # Extract from all shapes
    for shape in slide.shapes:
        # Text frames (text boxes, titles, content)
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                texts.append(text)

        # Tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        row_texts.append(text)
                if row_texts:
                    texts.append(" | ".join(row_texts))

        # Groups (nested shapes)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if hasattr(sub_shape, "text") and sub_shape.text:
                    text = sub_shape.text.strip()
                    if text:
                        texts.append(text)

    # Extract from notes
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            texts.append(f"Speaker Notes: {notes_text}")

    return {
        'slide_number': slide_number,
        'content': '\n'.join(texts),
        'text_count': len(texts)
    }

def extract_text_from_pptx(file_path):
    """Extract text from all slides in a PowerPoint file"""
    try:
        # Load presentation
        prs = Presentation(file_path)

        # Extract metadata
        metadata = {
            'title': prs.core_properties.title or '',
            'author': prs.core_properties.author or '',
            'subject': prs.core_properties.subject or '',
            'created': str(prs.core_properties.created) if prs.core_properties.created else '',
            'modified': str(prs.core_properties.modified) if prs.core_properties.modified else '',
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        }

        # Extract text from each slide
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_data = extract_text_from_slide(slide, idx)
            slides.append(slide_data)

        # Combine all text with slide markers
        full_text_parts = []
        for s in slides:
            if s['content']:
                full_text_parts.append(f"=== Slide {s['slide_number']} ===\n{s['content']}")

        full_text = '\n\n'.join(full_text_parts)

        return {
            'success': True,
            'metadata': metadata,
            'slides': slides,
            'full_text': full_text,
            'total_slides': len(prs.slides),
            'slides_with_text': len([s for s in slides if s['content']]),
            'total_characters': len(full_text)
        }

    except FileNotFoundError:
        return {
            'success': False,
            'error': f'File not found: {file_path}'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'Error extracting PPTX: {str(e)}'
        }

if __name__ == '__main__':
    if len(sys.argv) != 2:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python extract_pptx.py <file_path>'
        }))
        sys.exit(1)

    file_path = sys.argv[1]
    result = extract_text_from_pptx(file_path)
    print(json.dumps(result, indent=2, ensure_ascii=False))
