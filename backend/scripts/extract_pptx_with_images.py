#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract text AND images from PowerPoint (.pptx) files
Supports text extraction from slides, tables, notes, and nested shapes
Also extracts embedded images from each slide
"""

import sys
import json
import io
import os
import base64
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Set UTF-8 encoding for stdout on Windows
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

def extract_images_from_slide(slide, slide_number, output_dir):
    """Extract all images from a single slide"""
    images = []

    for shape_idx, shape in enumerate(slide.shapes):
        # Check if shape has an image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # Get the image data
                image = shape.image
                image_bytes = image.blob

                # Get the content type (e.g., 'image/jpeg', 'image/png')
                content_type = image.content_type
                ext = content_type.split('/')[-1] if '/' in content_type else 'png'

                # Generate filename
                filename = f"slide_{slide_number}_img_{shape_idx + 1}.{ext}"
                filepath = os.path.join(output_dir, filename)

                # Save the image
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)

                # Also create a base64 version for immediate use
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')

                images.append({
                    'filename': filename,
                    'path': filepath,
                    'content_type': content_type,
                    'size': len(image_bytes),
                    'base64': f"data:{content_type};base64,{image_base64[:100]}..."  # Truncated for logging
                })

            except Exception as e:
                print(f"Warning: Failed to extract image from slide {slide_number}, shape {shape_idx}: {e}", file=sys.stderr)

    return images

def extract_text_from_slide(slide, slide_number):
    """Extract all text from a single slide including tables and notes"""
    texts = []

    # Extract from all shapes
    for shape in slide.shapes:
        # Text frames (text boxes, titles, content)
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                texts.append(text)

        # Tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        row_texts.append(text)
                if row_texts:
                    texts.append(" | ".join(row_texts))

        # Groups (nested shapes)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if hasattr(sub_shape, "text") and sub_shape.text:
                    text = sub_shape.text.strip()
                    if text:
                        texts.append(text)

    # Extract from notes
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            texts.append(f"Speaker Notes: {notes_text}")

    return {
        'slide_number': slide_number,
        'content': '\n'.join(texts),
        'text_count': len(texts)
    }

def extract_pptx_data(file_path, output_dir=None):
    """Extract text and images from all slides in a PowerPoint file"""
    try:
        # Load presentation
        prs = Presentation(file_path)

        # Create output directory for images if specified
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        # Extract metadata
        metadata = {
            'title': prs.core_properties.title or '',
            'author': prs.core_properties.author or '',
            'subject': prs.core_properties.subject or '',
            'created': str(prs.core_properties.created) if prs.core_properties.created else '',
            'modified': str(prs.core_properties.modified) if prs.core_properties.modified else '',
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        }

        # Extract text and images from each slide
        slides = []
        all_images = []

        for idx, slide in enumerate(prs.slides, start=1):
            # Extract text
            slide_data = extract_text_from_slide(slide, idx)

            # Extract images if output directory is provided
            if output_dir:
                images = extract_images_from_slide(slide, idx, output_dir)
                slide_data['images'] = images
                all_images.extend(images)

            slides.append(slide_data)

        # Combine all text with slide markers
        full_text_parts = []
        for s in slides:
            if s['content']:
                full_text_parts.append(f"=== Slide {s['slide_number']} ===\n{s['content']}")

        full_text = '\n\n'.join(full_text_parts)

        return {
            'success': True,
            'metadata': metadata,
            'slides': slides,
            'full_text': full_text,
            'total_slides': len(prs.slides),
            'slides_with_text': len([s for s in slides if s['content']]),
            'total_characters': len(full_text),
            'total_images': len(all_images),
            'images': all_images
        }

    except FileNotFoundError:
        return {
            'success': False,
            'error': f'File not found: {file_path}'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'Error extracting PPTX: {str(e)}'
        }

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python extract_pptx_with_images.py <file_path> [output_dir]'
        }))
        sys.exit(1)

    file_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None

    result = extract_pptx_data(file_path, output_dir)
    print(json.dumps(result, indent=2, ensure_ascii=False))
